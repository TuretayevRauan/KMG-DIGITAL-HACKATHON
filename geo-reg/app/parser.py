import os
import io
import json
import glob
import base64
import hashlib
import csv as csv_module
import concurrent.futures
import fitz
import pdfplumber
import httpx
from docx import Document
from dotenv import load_dotenv

load_dotenv()

OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
OCR_BATCH_SIZE = int(os.getenv("OCR_BATCH_SIZE", "5"))

# OCR пробуем по списку vision-моделей с fallback (как answerer по тексту):
# если модель недоступна на ключе (404) или перегружена (429) — берём следующую.
# Сначала бесплатные vision-модели, затем платный надёжный fallback.
# Можно переопределить через env OCR_MODELS="model1,model2,...".
DEFAULT_OCR_MODELS = [
    "qwen/qwen3-vl-8b-instruct",             # дёшево, специализирована на OCR
    "qwen/qwen3-vl-30b-a3b-instruct",        # точнее, всё ещё дёшево
    "nvidia/nemotron-nano-12b-v2-vl:free",   # бесплатная VL-модель
    "google/gemma-4-31b-it:free",            # бесплатная, мультиязычная (open-weight)
    "google/gemma-4-26b-a4b-it:free",        # бесплатный fallback (open-weight)
]
_env_models = os.getenv("OCR_MODELS") or os.getenv("OCR_MODEL")
OCR_MODELS = (
    [m.strip() for m in _env_models.split(",") if m.strip()]
    if _env_models else DEFAULT_OCR_MODELS
)

# Минимальная длина текстового слоя, чтобы считать страницу "текстовой".
# Если на странице символов меньше — считаем её сканом и отправляем в OCR
# (страницы, где есть только подпись/штамп, но основное содержимое — картинка).
MIN_TEXT_LEN = int(os.getenv("MIN_TEXT_LEN", "20"))

# ---------- Локальный OCR (Tesseract) — основной путь, в разы быстрее API ----------
# Рендерим страницы PyMuPDF и распознаём локально через Tesseract в несколько
# потоков. API-vision остаётся как fallback. Результат кэшируется по хэшу файла,
# чтобы повторная загрузка того же PDF не запускала OCR заново.
USE_LOCAL_OCR = os.getenv("USE_LOCAL_OCR", "1") == "1"
OCR_DPI       = int(os.getenv("OCR_DPI", "200"))        # 200 dpi — баланс скорость/качество
OCR_LANG      = os.getenv("OCR_LANG", "rus+eng")        # нужен пакет языка rus
OCR_WORKERS   = int(os.getenv("OCR_WORKERS", str(min(8, (os.cpu_count() or 4)))))
OCR_CACHE_DIR = os.getenv("OCR_CACHE_DIR", ".ocr_cache")
# Путь к tesseract.exe. Если не в PATH — задаётся через env TESSERACT_CMD,
# либо ищется автоматически в типичных местах установки на Windows.
TESSERACT_CMD = os.getenv("TESSERACT_CMD", "").strip()


def _find_tesseract() -> str | None:
    """Возвращает путь к tesseract(.exe): env -> PATH -> типичные места установки."""
    import shutil
    if TESSERACT_CMD and os.path.isfile(TESSERACT_CMD):
        return TESSERACT_CMD
    found = shutil.which("tesseract")
    if found:
        return found
    candidates = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    ]
    for base in (os.getenv("LOCALAPPDATA"), os.getenv("USERPROFILE"), os.getenv("ProgramFiles")):
        if base:
            candidates.append(os.path.join(base, "Tesseract-OCR", "tesseract.exe"))
            candidates.append(os.path.join(base, "Programs", "Tesseract-OCR", "tesseract.exe"))
    for c in candidates:
        if c and os.path.isfile(c):
            return c
    # последний шанс — глубокий поиск по Program Files (медленно, но один раз)
    for base in (os.getenv("ProgramFiles"), os.getenv("ProgramFiles(x86)"), os.getenv("LOCALAPPDATA")):
        if base and os.path.isdir(base):
            hits = glob.glob(os.path.join(base, "**", "tesseract.exe"), recursive=True)
            if hits:
                return hits[0]
    return None


def _local_ocr_available() -> bool:
    try:
        import pytesseract
        from PIL import Image  # noqa: F401
        exe = _find_tesseract()
        if exe:
            pytesseract.pytesseract.tesseract_cmd = exe
        pytesseract.get_tesseract_version()
        return True
    except Exception as e:
        print(f"  (local OCR недоступен: {e})")
        return False


def _file_hash(file_path: str) -> str:
    h = hashlib.md5()
    with open(file_path, "rb") as f:
        for blk in iter(lambda: f.read(1 << 20), b""):
            h.update(blk)
    return h.hexdigest()


def _cache_path(file_hash: str) -> str:
    os.makedirs(OCR_CACHE_DIR, exist_ok=True)
    return os.path.join(OCR_CACHE_DIR, f"{file_hash}_{OCR_DPI}_{OCR_LANG}.json")


def _load_ocr_cache(file_hash: str) -> dict[int, str]:
    p = _cache_path(file_hash)
    if os.path.exists(p):
        try:
            with open(p, encoding="utf-8") as f:
                return {int(k): v for k, v in json.load(f).items()}
        except Exception:
            return {}
    return {}


def _save_ocr_cache(file_hash: str, data: dict[int, str]) -> None:
    try:
        with open(_cache_path(file_hash), "w", encoding="utf-8") as f:
            json.dump({str(k): v for k, v in data.items()}, f, ensure_ascii=False)
    except Exception as e:
        print(f"  (cache) не сохранил: {e}")


def local_ocr_pages(file_path: str, page_indices: list[int]) -> dict[int, str]:
    """OCR указанных страниц (0-based) локально через Tesseract в OCR_WORKERS потоков.
    Возвращает {page_index: text}. Рендер батчами, чтобы не держать всё в памяти."""
    import pytesseract
    from PIL import Image

    results: dict[int, str] = {}
    doc = fitz.open(file_path)
    mat = fitz.Matrix(OCR_DPI / 72, OCR_DPI / 72)
    batch = max(OCR_WORKERS * 4, OCR_WORKERS)
    try:
        for i in range(0, len(page_indices), batch):
            sub = page_indices[i: i + batch]
            pngs = {idx: doc[idx].get_pixmap(matrix=mat).tobytes("png") for idx in sub}

            def _ocr_one(idx: int):
                try:
                    img = Image.open(io.BytesIO(pngs[idx]))
                    return idx, pytesseract.image_to_string(img, lang=OCR_LANG).strip()
                except Exception as e:
                    print(f"  ⚠ local OCR стр.{idx+1}: {e}")
                    return idx, ""

            with concurrent.futures.ThreadPoolExecutor(max_workers=OCR_WORKERS) as ex:
                for idx, txt in ex.map(_ocr_one, sub):
                    results[idx] = txt
            print(f"  Local OCR: {min(i + batch, len(page_indices))}/{len(page_indices)} страниц")
    finally:
        doc.close()
    return results

# Какие форматы поддерживаем
PDF_EXT = {".pdf"}
DOCX_EXT = {".docx"}
TEXT_EXT = {".txt", ".md", ".markdown", ".log"}
CSV_EXT = {".csv", ".tsv"}
XLSX_EXT = {".xlsx", ".xlsm", ".xls"}
PPTX_EXT = {".pptx"}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp", ".gif"}

SUPPORTED_EXT = (
    PDF_EXT | DOCX_EXT | TEXT_EXT | CSV_EXT | XLSX_EXT | PPTX_EXT | IMAGE_EXT
)


def parse_document(file_path: str, doc_id: str) -> list[dict]:
    """Читает документ ЛЮБОГО поддерживаемого формата и возвращает чанки.
    PDF: текстовый слой + таблицы + OCR сканов. Картинки/сканы — OCR.
    Office/текст/таблицы — напрямую."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext in PDF_EXT:
        return parse_pdf(file_path, doc_id)
    if ext in DOCX_EXT:
        return parse_docx(file_path, doc_id)
    if ext in TEXT_EXT:
        return parse_plain_text(file_path, doc_id)
    if ext in CSV_EXT:
        return parse_csv(file_path, doc_id)
    if ext in XLSX_EXT:
        return parse_xlsx(file_path, doc_id)
    if ext in PPTX_EXT:
        return parse_pptx(file_path, doc_id)
    if ext in IMAGE_EXT:
        return parse_image(file_path, doc_id)
    raise ValueError(
        f"Формат {ext} не поддерживается. Доступно: {', '.join(sorted(SUPPORTED_EXT))}"
    )


# ---------- OCR ----------

def pdf_page_to_base64(page, dpi: int = 150) -> str:
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat)
    return base64.b64encode(pix.tobytes("png")).decode("utf-8")


def ocr_pages_with_openrouter(images_b64: list[str], page_numbers: list[int]) -> dict[int, str]:
    """OCR страниц/картинок через OpenRouter vision-модель (open-weight, см.
    OCR_MODELS). Используется как fallback после локального Tesseract.
    Тот же OPENROUTER_API_KEY, что и в answerer."""
    if not OPENROUTER_API_KEY:
        raise ValueError("OPENROUTER_API_KEY не задан в .env файле!")

    content = []
    for img_b64, page_num in zip(images_b64, page_numbers):
        content.append({"type": "text", "text": f"Страница {page_num}:"})
        content.append({
            "type": "image_url",
            "image_url": {"url": f"data:image/png;base64,{img_b64}"},
        })

    content.append({
        "type": "text",
        "text": f"""Извлеки ВЕСЬ текст с каждой страницы/картинки выше как можно точнее.
Верни ТОЛЬКО JSON без пояснений:
{{{", ".join(f'"{p}": "текст страницы {p}"' for p in page_numbers)}}}

Правила:
- Сохраняй структуру: заголовки, абзацы, списки.
- Таблицы передавай в виде строк, ячейки через " | ".
- Числа, единицы измерения, индексы скважин и пластов переноси точно.
- Ничего не придумывай: если текст нечитаем — оставь пустую строку.""",
    })

    # Перебираем vision-модели по очереди: первая, что ответит — выигрывает.
    last_err = None
    for model in OCR_MODELS:
        try:
            print(f"  OCR страниц {page_numbers} через {model}...")
            response = httpx.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                    "Content-Type": "application/json",
                    "HTTP-Referer": "http://localhost:8000",
                    "X-Title": "GeoAI OCR",
                },
                json={
                    "model": model,
                    "max_tokens": 4096,
                    "messages": [{"role": "user", "content": content}],
                },
                timeout=120.0,
            )
            if response.status_code in (404, 429, 402, 403, 503):
                print(f"  ⚠ {model} недоступна ({response.status_code}), пробую следующую...")
                last_err = f"{model}: HTTP {response.status_code}"
                continue
            response.raise_for_status()
            data = response.json()
            if "choices" not in data or not data["choices"]:
                print(f"  ⚠ {model} вернула пустой ответ ({data.get('error')}), следующая...")
                last_err = f"{model}: {data.get('error')}"
                continue
            raw = data["choices"][0]["message"]["content"].strip()
            if raw.startswith("```"):
                raw = raw.split("\n", 1)[1].rsplit("```", 1)[0].strip()
            try:
                parsed = json.loads(raw)
                result = {int(k): v for k, v in parsed.items()}
            except (json.JSONDecodeError, ValueError):
                # Модель вернула не-JSON (просто текст) — отдаём как одну страницу.
                if len(page_numbers) == 1:
                    result = {page_numbers[0]: raw}
                else:
                    print(f"  ⚠ {model}: ответ не JSON, следующая...")
                    last_err = f"{model}: non-JSON"
                    continue
            if any(v and v.strip() for v in result.values()):
                print(f"  ✓ OCR удался через {model}")
                return result
            print(f"  ⚠ {model}: пустой текст, следующая...")
            last_err = f"{model}: empty"
        except Exception as e:
            print(f"  ⚠ Ошибка OCR ({model}): {e}, следующая...")
            last_err = f"{model}: {e}"
            continue

    print(f"  ✗ Все OCR-модели не сработали. Последняя ошибка: {last_err}")
    return {p: "" for p in page_numbers}


# ---------- PDF ----------

def extract_all_tables(file_path: str, total_pages: int) -> dict[int, str]:
    """Извлекает таблицы со всех страниц за одно открытие файла."""
    tables_by_page: dict[int, str] = {}
    try:
        with pdfplumber.open(file_path) as pdf:
            for page_num in range(min(total_pages, len(pdf.pages))):
                tables = pdf.pages[page_num].extract_tables()
                text = ""
                for table in tables:
                    if table:
                        rows = [
                            " | ".join(str(cell) if cell else "" for cell in row)
                            for row in table
                        ]
                        text += "\n".join(rows) + "\n\n"
                if text:
                    tables_by_page[page_num] = text
    except Exception as e:
        print(f"Ошибка при извлечении таблиц: {e}")
    return tables_by_page


def parse_pdf(file_path: str, doc_id: str) -> list[dict]:
    chunks = []
    doc = fitz.open(file_path)
    total_pages = len(doc)
    print(f"Обрабатываем PDF: {total_pages} страниц")

    # Текстовые страницы vs сканы (по порогу длины текстового слоя)
    text_pages: dict[int, str] = {}
    scan_page_indices: list[int] = []
    for page_num in range(total_pages):
        text = doc[page_num].get_text().strip()
        if len(text) >= MIN_TEXT_LEN:
            text_pages[page_num] = text
        else:
            scan_page_indices.append(page_num)
    print(f"  Текстовых: {len(text_pages)}, сканов: {len(scan_page_indices)}")

    doc.close()

    # OCR сканов: кэш -> локальный Tesseract (быстро, параллельно) -> API-vision (fallback)
    ocr_results: dict[int, str] = {}   # ключ — 1-based номер страницы (как раньше)
    if scan_page_indices:
        file_hash = _file_hash(file_path)
        cache = _load_ocr_cache(file_hash)

        cached_idx = [i for i in scan_page_indices if cache.get(i, "").strip()]
        todo_idx   = [i for i in scan_page_indices if not cache.get(i, "").strip()]
        for i in cached_idx:
            ocr_results[i + 1] = cache[i]
        if cached_idx:
            print(f"  Из кэша OCR: {len(cached_idx)} страниц")

        if todo_idx and USE_LOCAL_OCR and _local_ocr_available():
            print(f"  Локальный Tesseract OCR: {len(todo_idx)} стр., "
                  f"{OCR_WORKERS} потоков, {OCR_DPI} dpi, lang='{OCR_LANG}'...")
            for i, txt in local_ocr_pages(file_path, todo_idx).items():
                if txt.strip():
                    ocr_results[i + 1] = txt
                    cache[i] = txt
            todo_idx = [i for i in todo_idx if not ocr_results.get(i + 1, "").strip()]
        elif todo_idx:
            print("  ⚠ Локальный Tesseract недоступен — иду через API OCR. "
                  "Для скорости установи: pip install pytesseract pillow + Tesseract (rus).")

        # API-vision только для того, что осталось (локально не распозналось)
        if todo_idx:
            print(f"  API OCR (fallback): {len(todo_idx)} стр., батчами по {OCR_BATCH_SIZE}...")
            doc2 = fitz.open(file_path)
            try:
                for i in range(0, len(todo_idx), OCR_BATCH_SIZE):
                    batch_indices = todo_idx[i: i + OCR_BATCH_SIZE]
                    images_b64 = [pdf_page_to_base64(doc2[idx]) for idx in batch_indices]
                    page_numbers = [idx + 1 for idx in batch_indices]
                    for pnum, txt in ocr_pages_with_openrouter(images_b64, page_numbers).items():
                        if txt and txt.strip():
                            ocr_results[pnum] = txt
                            cache[pnum - 1] = txt
                    print(f"  API OCR: {min(i + OCR_BATCH_SIZE, len(todo_idx))}/{len(todo_idx)}")
            finally:
                doc2.close()

        _save_ocr_cache(file_hash, cache)

    tables_by_page = extract_all_tables(file_path, total_pages)

    for page_num in range(total_pages):
        if page_num in text_pages:
            full_text = text_pages[page_num]
            if page_num in tables_by_page:
                full_text += "\n\n[ТАБЛИЦЫ]\n" + tables_by_page[page_num]
        else:
            full_text = ocr_results.get(page_num + 1, "")

        if not full_text.strip():
            continue

        for i, chunk_text in enumerate(split_into_chunks(full_text)):
            chunks.append(_chunk(doc_id, page_num + 1, i, chunk_text))

    # Описания графики (ТЗ 6.2) — чанки figure_caption (если ENABLE_VISION=1)
    try:
        from vision import extract_figure_captions
        chunks.extend(extract_figure_captions(file_path, doc_id))
    except Exception as e:
        print(f"  (предупреждение) vision не сработал: {e}")

    print(f"PDF распарсен: {len(chunks)} чанков из {total_pages} страниц")
    return chunks


# ---------- DOCX ----------

def parse_docx(file_path: str, doc_id: str) -> list[dict]:
    doc = Document(file_path)
    full_text = ""
    for para in doc.paragraphs:
        if para.text.strip():
            full_text += para.text + "\n"
    for table in doc.tables:
        for row in table.rows:
            full_text += " | ".join(cell.text for cell in row.cells) + "\n"
    return _chunks_single_page(full_text, doc_id)


# ---------- Plain text / Markdown ----------

def parse_plain_text(file_path: str, doc_id: str) -> list[dict]:
    text = _read_text_any_encoding(file_path)
    return _chunks_single_page(text, doc_id)


def _read_text_any_encoding(file_path: str) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    with open(file_path, "rb") as f:
        return f.read().decode("utf-8", errors="replace")


# ---------- CSV / TSV ----------

def parse_csv(file_path: str, doc_id: str) -> list[dict]:
    ext = os.path.splitext(file_path)[1].lower()
    delimiter = "\t" if ext == ".tsv" else ","
    raw = _read_text_any_encoding(file_path)
    lines = []
    reader = csv_module.reader(raw.splitlines(), delimiter=delimiter)
    for row in reader:
        if any(cell.strip() for cell in row):
            lines.append(" | ".join(cell.strip() for cell in row))
    return _chunks_single_page("\n".join(lines), doc_id)


# ---------- XLSX / XLS ----------

def parse_xlsx(file_path: str, doc_id: str) -> list[dict]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError(
            "Для чтения Excel установи: pip install openpyxl"
        )
    wb = load_workbook(file_path, read_only=True, data_only=True)
    chunks: list[dict] = []
    sheet_count = len(wb.worksheets)
    for sheet_idx, ws in enumerate(wb.worksheets, start=1):
        lines = [f"[Лист: {ws.title}]"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                lines.append(" | ".join(cells))
        sheet_text = "\n".join(lines)
        if sheet_text.strip():
            for i, chunk_text in enumerate(split_into_chunks(sheet_text)):
                chunks.append(_chunk(doc_id, sheet_idx, i, chunk_text))
    wb.close()
    print(f"Excel распарсен: {len(chunks)} чанков, листов: {sheet_count}")
    return chunks


# ---------- PPTX ----------

def parse_pptx(file_path: str, doc_id: str) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "Для чтения PowerPoint установи: pip install python-pptx"
        )
    prs = Presentation(file_path)
    chunks: list[dict] = []
    slide_count = len(prs.slides._sldIdLst)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
        slide_text = "\n".join(parts)
        if slide_text.strip():
            for i, chunk_text in enumerate(split_into_chunks(slide_text)):
                chunks.append(_chunk(doc_id, slide_idx, i, chunk_text))
    print(f"PPTX распарсен: {len(chunks)} чанков, слайдов: {slide_count}")
    return chunks


# ---------- Images (OCR) ----------

def parse_image(file_path: str, doc_id: str) -> list[dict]:
    # Нормализуем любую картинку (jpg/tiff/bmp/webp/…) в PNG через PyMuPDF,
    # чтобы vision-модель точно её приняла.
    try:
        img_doc = fitz.open(file_path)
        pix = img_doc[0].get_pixmap(matrix=fitz.Matrix(2, 2))
        img_b64 = base64.b64encode(pix.tobytes("png")).decode("utf-8")
        img_doc.close()
    except Exception:
        with open(file_path, "rb") as f:
            img_b64 = base64.b64encode(f.read()).decode("utf-8")
    ocr = ocr_pages_with_openrouter([img_b64], [1])
    text = ocr.get(1, "")
    return _chunks_single_page(text, doc_id)


# ---------- Общие хелперы ----------

def _chunk(doc_id: str, page: int, idx: int, text: str) -> dict:
    return {
        "chunk_id": f"{doc_id}_p{page}_c{idx}",
        "doc_id": doc_id,
        "page": page,
        "text": text,
        "type": "text",
    }


def _chunks_single_page(text: str, doc_id: str) -> list[dict]:
    return [
        _chunk(doc_id, 1, i, chunk)
        for i, chunk in enumerate(split_into_chunks(text))
    ]


CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "1400"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "150"))


def split_into_chunks(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    if len(text) <= chunk_size:
        return [text.strip()] if text.strip() else []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if end < len(text):
            bp = max(chunk.rfind("."), chunk.rfind("\n"))
            if bp > chunk_size // 2:
                chunk = text[start: start + bp + 1]
                end = start + bp + 1
        chunks.append(chunk.strip())
        start = end - overlap
    return [c for c in chunks if c]
